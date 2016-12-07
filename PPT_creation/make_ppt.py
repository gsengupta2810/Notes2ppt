from pptx import Presentation
import text_preparation_for_ppt as tpp
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor

MAX_NUMBER_BULLETS_PER_SLIDE=5

def add_slide(prs,SLD_LAYOUT_TITLE_AND_CONTENT = 1):
	slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
	slide = prs.slides.add_slide(slide_layout)
	return slide

def add_text_to_slide(slide,topic,data):
	shapes = slide.shapes
	slide_placeholders = slide.placeholders
	for i in xrange(0,len(slide_placeholders)):
		if (slide_placeholders[i].name).split()[0]=="Content":
			body_shape = shapes.placeholders[1]
			tf = body_shape.text_frame
			p=tf.add_paragraph()
			p.text=data
			p.level=0
			p.font.size=Pt(18)
			p.font.color.rgb = RGBColor(12, 34, 56)
			tf.margin_top=Inches(0.08)
			tf.margin_bottom=Inches(0.08)
			tf.word_wrap = True
		elif (slide_placeholders[i].name).split()[0]=="Title":
			title_shape = shapes.title
			title_shape.text = topic
			title_shape.margin_bottom=Pt(0)
	return slide

if __name__ == '__main__':
	prs = Presentation()
	file=open('./test_extracted_data_for_ppt.txt','r')
	topics,topic_data=tpp.split_sentences_nltk(file)
	for i in xrange(0,len(topic_data)):
		slide=add_slide(prs,1)
		content=topic_data[i].split('\n')
		counter=0
		for j in xrange(0,len(content)):
			if counter<MAX_NUMBER_BULLETS_PER_SLIDE:
				add_text_to_slide(slide,topics[i],content[j])
				counter+=1
			else:
				slide=add_slide(prs,1)
				add_text_to_slide(slide,topics[i],content[j])
				counter=0
	prs.save('test.pptx')


